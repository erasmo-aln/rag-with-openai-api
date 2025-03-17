import os
import glob
import fitz
import chromadb
import dotenv
from docx import Document
from pptx import Presentation
from langchain_openai import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter


# Load API key
dotenv.load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

# Initialize embedding model
embedding_model = OpenAIEmbeddings(model="text-embedding-ada-002", openai_api_key=OPENAI_API_KEY)

# Initialize ChromaDB
chroma_client = chromadb.PersistentClient(path="./chroma_db")
collection = chroma_client.get_or_create_collection(name="documents")

# Function to extract text from PDFs
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = "\n".join(page.get_text("text") for page in doc)
    return text.strip()


# Function to extract text from DOCX
def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs]).strip()
    return text


# Function to extract text from PPTX (1 slide = 1 chunk)
def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides.append("\n".join(slide_text).strip())
    return slides


# Function to check if a file is already stored in ChromaDB
def is_file_stored(filename):
    return filename in collection.get(include=["metadatas"]).get("metadatas", [])


# Function to process and store documents
def ingest_documents():
    file_types = ["files/*.pdf", "files/*.docx", "files/*.pptx"]

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=100)

    for file_type in file_types:
        for file_path in glob.glob(file_type):
            file_name = os.path.basename(file_path)
            if is_file_stored(file_name):
                print(f"Skipping {file_name}, already stored.")
                continue

            print(f"Processing {file_name}...")

            # Extract text based on file type
            if file_path.endswith(".pdf"):
                text = extract_text_from_pdf(file_path)
                chunks = text_splitter.split_text(text)

            elif file_path.endswith(".docx"):
                text = extract_text_from_docx(file_path)
                chunks = text_splitter.split_text(text)

            elif file_path.endswith(".pptx"):
                chunks = extract_text_from_pptx(file_path)

            else:
                continue

            # Generate embeddings for each chunk
            embeddings = embedding_model.embed_documents(chunks)

            # Store chunks in ChromaDB
            for i, emb in enumerate(embeddings):
                collection.add(
                    ids=[f"{file_name}_{i}"],
                    documents=[chunks[i]],
                    embeddings=[emb],
                    metadatas=[{"filename": file_name}]  # ✅ Ensures document name is stored
                )


            print(f"Stored {file_name} in ChromaDB.")



if __name__ == "__main__":
    ingest_documents()
